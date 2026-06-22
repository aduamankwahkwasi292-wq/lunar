"""
Document text extraction for Lunar AI.

Handles the common case the user cares about most: slides that are really just
images. Native text is extracted where present, and OCR (RapidOCR, ONNX, fully
bundled — no Tesseract/Poppler system install needed) is used for:

  * standalone image files,
  * image-only / scanned PDF pages (rasterized with PyMuPDF), and
  * pictures embedded inside PowerPoint slides.
"""

import io
import os
import threading

# Minimum characters a PDF page / image must yield before we trust the native
# text layer; below this we fall back to OCR.
_MIN_NATIVE_CHARS = 16

_ocr_engine = None
_ocr_lock = threading.Lock()
_ocr_available = True


def extract_text_from_file(file_path: str) -> str:
    """Extract text content from various file types."""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == '.pdf':
            return _extract_pdf(file_path)
        elif ext in ('.docx', '.doc'):
            return _extract_docx(file_path)
        elif ext in ('.pptx', '.ppt'):
            return _extract_pptx(file_path)
        elif ext == '.xlsx':
            return _extract_xlsx(file_path)
        elif ext == '.csv':
            return _extract_csv(file_path)
        elif ext in ('.txt', '.md', '.log', '.py', '.js', '.html', '.css', '.json',
                      '.xml', '.yaml', '.yml', '.ini', '.cfg', '.env', '.sh', '.bat',
                      '.c', '.cpp', '.h', '.java', '.rb', '.go', '.rs', '.ts', '.tsx',
                      '.jsx', '.sql', '.r', '.m', '.swift', '.kt', '.scala'):
            return _extract_text(file_path)
        elif ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif', '.webp'):
            return _extract_image(file_path)
        else:
            return _extract_text(file_path)
    except Exception as e:
        return f"[Error extracting text from {os.path.basename(file_path)}: {str(e)}]"


# ======================================================================
# OCR engine (RapidOCR / ONNX) — lazy, bundled, offline
# ======================================================================
def _get_ocr():
    """Lazily build the RapidOCR engine; returns None if it can't be loaded."""
    global _ocr_engine, _ocr_available
    if _ocr_engine is not None or not _ocr_available:
        return _ocr_engine
    with _ocr_lock:
        if _ocr_engine is None and _ocr_available:
            try:
                from rapidocr_onnxruntime import RapidOCR
                _ocr_engine = RapidOCR()
            except Exception:
                _ocr_available = False
                _ocr_engine = None
    return _ocr_engine


def _ocr_image_obj(image) -> str:
    """Run OCR on a PIL image (or path). Returns recognized text, or ''."""
    engine = _get_ocr()
    if engine is None:
        return _ocr_with_tesseract(image)
    try:
        import numpy as np
        from PIL import Image

        if not isinstance(image, Image.Image):
            image = Image.open(image)
        arr = np.array(image.convert("RGB"))
        result, _elapse = engine(arr)
        if not result:
            return ""
        # result rows are [box, text, score]
        return "\n".join(row[1] for row in result if len(row) > 1 and row[1]).strip()
    except Exception:
        return _ocr_with_tesseract(image)


def _ocr_with_tesseract(image) -> str:
    """Optional fallback if RapidOCR is unavailable but Tesseract is installed."""
    try:
        import pytesseract
        from PIL import Image
        if not isinstance(image, Image.Image):
            image = Image.open(image)
        return (pytesseract.image_to_string(image) or "").strip()
    except Exception:
        return ""


# ======================================================================
# Per-format extractors
# ======================================================================
def _extract_pdf(file_path: str) -> str:
    """Extract PDF text; OCR any page whose native text layer is empty/sparse."""
    text_parts = []
    native_by_page = []
    try:
        import PyPDF2
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                native_by_page.append((page.extract_text() or "").strip())
    except ImportError:
        native_by_page = []
    except Exception:
        native_by_page = []

    # Rasterize + OCR pages that have little or no native text.
    rasterizer = _open_pdf_for_render(file_path)
    for idx, native in enumerate(native_by_page or []):
        if len(native) >= _MIN_NATIVE_CHARS:
            text_parts.append(native)
        elif rasterizer is not None:
            ocr_text = _ocr_pdf_page(rasterizer, idx)
            text_parts.append(ocr_text or native)
        else:
            text_parts.append(native)

    # PyPDF2 unavailable but we can still render: OCR the whole document.
    if not native_by_page and rasterizer is not None:
        for idx in range(rasterizer.page_count):
            text_parts.append(_ocr_pdf_page(rasterizer, idx))

    if rasterizer is not None:
        try:
            rasterizer.close()
        except Exception:
            pass

    combined = "\n".join(p for p in text_parts if p).strip()
    if combined:
        return combined
    if not native_by_page and rasterizer is None:
        return "[Could not read PDF. Install PyPDF2 and PyMuPDF for text + OCR support.]"
    return combined


def _open_pdf_for_render(file_path: str):
    try:
        import fitz  # PyMuPDF
        return fitz.open(file_path)
    except Exception:
        return None


def _ocr_pdf_page(doc, page_index: int, dpi: int = 200) -> str:
    """Render a single PDF page to an image and OCR it."""
    try:
        from PIL import Image
        page = doc.load_page(page_index)
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        return _ocr_image_obj(img)
    except Exception:
        return ""


def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
    except ImportError:
        return "[python-docx not installed. Install with: pip install python-docx]"


def _extract_pptx(file_path: str) -> str:
    """Extract slide text boxes AND OCR any pictures embedded in the slides."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return "[python-pptx not installed. Install with: pip install python-pptx]"

    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # Native text in text boxes / placeholders.
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                text_parts.append(shape.text.strip())
            # OCR pictures (image-based slides).
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    from PIL import Image
                    blob = shape.image.blob
                    ocr_text = _ocr_image_obj(Image.open(io.BytesIO(blob)))
                    if ocr_text:
                        text_parts.append(ocr_text)
            except Exception:
                pass
    return '\n'.join(text_parts)


def _extract_xlsx(file_path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text_parts.append(row_text)
        return '\n'.join(text_parts)
    except ImportError:
        return "[openpyxl not installed. Install with: pip install openpyxl]"


def _extract_csv(file_path: str) -> str:
    import csv
    text_parts = []
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            text_parts.append(' '.join(row))
    return '\n'.join(text_parts)


def _extract_text(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _extract_image(file_path: str) -> str:
    text = _ocr_image_obj(file_path)
    return text if text else "[No text found in image]"


# ======================================================================
# Slide-by-slide extraction (for the classroom: one board lesson per slide)
# ======================================================================
import re as _re


def extract_slides(file_path: str) -> list:
    """Split a document into its real slides/pages.

    Returns a list of slide texts (one per PPTX slide / PDF page / image),
    or an auto-split into slide-sized sections for plain documents. Empty
    slides are dropped.
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext in ('.pptx', '.ppt'):
            slides = _pptx_slides(file_path)
        elif ext == '.pdf':
            slides = _pdf_slides(file_path)
        elif ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif', '.webp'):
            slides = [_extract_image(file_path)]
        else:
            slides = _autosplit(extract_text_from_file(file_path))
    except Exception as e:  # noqa: BLE001
        slides = _autosplit(f"[Error reading {os.path.basename(file_path)}: {e}]")

    cleaned = [s.strip() for s in slides if s and s.strip() and not s.strip().startswith("[")]
    return cleaned or [s.strip() for s in slides if s and s.strip()]


def _pptx_slides(file_path: str) -> list:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return _autosplit(_extract_pptx(file_path))

    prs = Presentation(file_path)
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    from PIL import Image
                    ocr = _ocr_image_obj(Image.open(io.BytesIO(shape.image.blob)))
                    if ocr:
                        parts.append(ocr)
            except Exception:
                pass
        slides.append("\n".join(parts))
    return slides


def _pdf_slides(file_path: str) -> list:
    native = []
    try:
        import PyPDF2
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            native = [(p.extract_text() or "").strip() for p in reader.pages]
    except Exception:
        native = []

    rasterizer = _open_pdf_for_render(file_path)
    pages = []
    page_count = len(native) if native else (rasterizer.page_count if rasterizer else 0)
    for idx in range(page_count):
        text = native[idx] if idx < len(native) else ""
        if len(text) < _MIN_NATIVE_CHARS and rasterizer is not None:
            text = _ocr_pdf_page(rasterizer, idx) or text
        pages.append(text)
    if rasterizer is not None:
        try:
            rasterizer.close()
        except Exception:
            pass
    return pages or _autosplit(_extract_pdf(file_path))


def _autosplit(text: str, target: int = 850) -> list:
    """Split plain text into slide-sized sections respecting natural breaks."""
    text = (text or "").strip()
    if not text:
        return []

    # Explicit page breaks (form feed) from PDF-style extraction.
    if "\f" in text:
        parts = [p.strip() for p in text.split("\f") if p.strip()]
        if len(parts) > 1:
            return parts

    # Explicit "Slide N" / "Page N" / "Section N" markers.
    marked = _re.split(r"\n(?=(?:slide|page|section|chapter|topic|unit)\s*\d+\b)", text, flags=_re.I)
    if len(marked) > 1:
        return [m.strip() for m in marked if m.strip()]

    # Otherwise pack paragraphs up to ~target chars.
    paras = [p.strip() for p in _re.split(r"\n\s*\n", text) if p.strip()] or [text]
    slides, buf = [], ""
    for p in paras:
        if buf and len(buf) + len(p) + 2 > target:
            slides.append(buf)
            buf = p
        else:
            buf = f"{buf}\n\n{p}".strip() if buf else p
    if buf:
        slides.append(buf)
    return slides


# ======================================================================
# Rich slide extraction — verbatim text + inline media (images / video) +
# links, as ordered "blocks" the classroom board renders in slide order.
# Each slide is:
#   {"title": str,
#    "blocks": [ {"type":"text","text":..}
#              | {"type":"image","src":<filename>,"alt":..}
#              | {"type":"video","src":<filename>} ],
#    "text":  "<verbatim text only>",   # what gets typed on the board
#    "ocr":   "<text recognised inside figures>",  # grounding only
#    "links": [ {"text":.., "url":..} ],
#    "examinable": bool }
# `src` is a bare filename inside media_dir; the web layer maps it to a URL.
# ======================================================================

# A picture covering more than this fraction of the whole slide/page is a
# full-slide scan: we OCR it to text rather than pasting the slide as an
# image (the user explicitly does not want the slide dumped as a picture).
_FULLSLIDE_IMG_FRAC = 0.7


def extract_rich_slides(file_path: str, media_dir: str) -> list:
    """Split a document into slides.

    PDFs and images become "board" slides: the page is rendered (lazily, on
    demand) as a chalk-on-transparent image so it sits on the green board with
    perfect math symbols / figures / number lines — native text is kept only for
    grounding. PPTX (no local renderer) keeps the verbatim text+figure blocks.
    """
    try:
        os.makedirs(media_dir, exist_ok=True)
    except Exception:
        pass
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext in ('.pptx', '.ppt'):
            slides = _pptx_rich(file_path, media_dir)
        elif ext == '.pdf':
            slides = _pdf_slides_meta(file_path)
        elif ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif', '.webp'):
            slides = _image_slide_meta(file_path)
        else:
            slides = _text_rich(extract_text_from_file(file_path))
    except Exception as e:  # noqa: BLE001
        slides = _text_rich(f"[Error reading {os.path.basename(file_path)}: {e}]")

    cleaned = []
    for s in slides:
        if (s.get("text") or "").strip() or s.get("blocks") or s.get("board"):
            s.setdefault("title", _first_line(s.get("text", "")))
            cleaned.append(s)
    return cleaned


# ----- board slides: render the real page recolored to chalk-on-board -----
def _page_links(page) -> list:
    links = []
    try:
        for ln in page.get_links():
            uri = ln.get("uri")
            if uri:
                links.append({"text": uri, "url": uri})
    except Exception:
        pass
    return _dedupe_links(links)


def _pdf_slides_meta(file_path: str) -> list:
    """One board slide per PDF page: native text (grounding) + lazy chalk render."""
    try:
        import fitz  # PyMuPDF
    except Exception:
        return _text_rich(_extract_pdf(file_path))
    doc = fitz.open(file_path)
    out = []
    try:
        for pi in range(doc.page_count):
            page = doc.load_page(pi)
            text = (page.get_text() or "").strip()
            # Blank page: nothing but maybe a page number, no images, no real
            # drawings. Cheap native check (no render) so we can tell the user.
            compact = _re.sub(r"\s+", "", text)
            page_num_only = bool(_re.fullmatch(r"\d{1,4}", compact or ""))
            has_content = len(compact) > 2 and not page_num_only
            try:
                n_draw = len(page.get_drawings())
                n_img = len(page.get_images())
            except Exception:
                n_draw, n_img = 0, 0
            blank = (not has_content) and n_img == 0 and n_draw <= 2
            out.append({
                "board": "pdf", "page": pi,
                "text": text, "ocr": "",
                "links": _page_links(page),
                "examinable": True,
                "blank": blank,
                "title": _first_line(text) or f"Slide {pi + 1}",
            })
    finally:
        try:
            doc.close()
        except Exception:
            pass
    return out


def _image_slide_meta(file_path: str) -> list:
    """A single image file becomes one board slide (recolored), OCR for grounding."""
    ocr = _ocr_image_obj(file_path) or ""
    return [{
        "board": "image", "page": 0,
        "text": ocr, "ocr": ocr, "links": [],
        "examinable": True,
        "title": _first_line(ocr) or "Slide",
    }]


import numpy as _np

# Precomputed alpha lookup: white paper (255) → 0 (transparent), ink (dark) →
# opaque, faint strokes boosted. A LUT gather is far cheaper than per-pixel math.
_ALPHA_LUT = _np.clip((255 - _np.arange(256, dtype=_np.float32) - 12.0) * 1.7, 0, 255).astype(_np.uint8)
_CHALK = (234, 240, 231)   # warm chalk white


def _chalk_from_rgb(arr):
    """Recolor an RGB(A) page raster to chalk-on-transparent RGBA.

    White paper → transparent (the board shows through); everything darker →
    warm chalk, opacity by how dark it is — the authentic chalkboard look. Pure
    integer / LUT work (no axis-2 reductions or float math) so a full page
    recolors in tens of ms. Returns (rgba_uint8_array, ink_coverage_fraction).
    """
    import numpy as np

    a = arr
    if a.ndim == 2:
        gray = a
    else:
        r = a[..., 0].astype(np.uint16)
        g = a[..., 1].astype(np.uint16)
        b = a[..., 2].astype(np.uint16)
        gray = ((r * 77 + g * 150 + b * 29) >> 8).astype(np.uint8)   # fast luminance
    alpha = _ALPHA_LUT[gray]
    h, w = gray.shape
    out = np.empty((h, w, 4), dtype=np.uint8)
    out[..., 0] = _CHALK[0]
    out[..., 1] = _CHALK[1]
    out[..., 2] = _CHALK[2]
    out[..., 3] = alpha
    coverage = float((alpha > 20).mean())
    return out, coverage


def _save_chalk_png(rgba, out_path: str):
    """Encode an RGBA chalk array to PNG quickly (low compression — cached anyway)."""
    from PIL import Image
    Image.fromarray(rgba, "RGBA").save(out_path, format="PNG", compress_level=1)


def render_pdf_page_chalk(file_path: str, page_index: int, out_path: str, dpi: int = 130) -> str:
    """Render a PDF page to a chalk-on-board PNG at out_path.

    Reads the rasterised pixels straight from the PyMuPDF pixmap (no intermediate
    PNG encode/decode), recolors, then writes one PNG — much faster than the
    previous double round-trip.
    """
    import fitz
    import numpy as np

    doc = fitz.open(file_path)
    try:
        page = doc.load_page(max(0, min(page_index, doc.page_count - 1)))
        pix = page.get_pixmap(dpi=dpi, alpha=False)
        arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
        if pix.n >= 4:
            arr = arr[..., :3]
        elif pix.n == 1:
            arr = np.repeat(arr, 3, axis=2)
        rgba, _ = _chalk_from_rgb(arr)
    finally:
        try:
            doc.close()
        except Exception:
            pass
    _save_chalk_png(rgba, out_path)
    return out_path


def pdf_page_words(file_path: str, page_index: int) -> list:
    """Word boxes for a PDF page as fractions of the page (for a selectable text layer).

    Returns a reading-ordered list of {x, y, w, h, t} where x/y/w/h are 0..1
    fractions of the page width/height, so the frontend can overlay an invisible,
    selectable <span> per word on top of the chalk image at any display size.
    """
    import fitz
    doc = fitz.open(file_path)
    try:
        page = doc.load_page(max(0, min(page_index, doc.page_count - 1)))
        pw = float(page.rect.width) or 1.0
        ph = float(page.rect.height) or 1.0
        # sort by block, line, word so selection follows reading order
        raw = page.get_text("words")
    finally:
        try:
            doc.close()
        except Exception:
            pass
    raw.sort(key=lambda w: (w[5], w[6], w[7]))
    out = []
    for x0, y0, x1, y1, text, *_rest in raw:
        t = (text or "").strip()
        if not t:
            continue
        out.append({
            "x": round(x0 / pw, 5), "y": round(y0 / ph, 5),
            "w": round((x1 - x0) / pw, 5), "h": round((y1 - y0) / ph, 5),
            "t": t,
        })
    return out


def pdf_page_lines(file_path: str, page_index: int) -> list:
    """Line boxes for a PDF page as page fractions, for a selectable text layer.

    Groups words into their text lines (one span per line) so the on-screen
    selection follows whole lines cleanly — dragging a sentence no longer bleeds
    into the line above or below the way per-word spans do. Returns reading-ordered
    {x, y, w, h, t} where t is the full line text and x/y/w/h are 0..1 fractions.
    """
    import fitz
    doc = fitz.open(file_path)
    try:
        page = doc.load_page(max(0, min(page_index, doc.page_count - 1)))
        pw = float(page.rect.width) or 1.0
        ph = float(page.rect.height) or 1.0
        raw = page.get_text("words")
    finally:
        try:
            doc.close()
        except Exception:
            pass
    # group by (block, line); each word tuple is (x0,y0,x1,y1,text,block,line,word)
    groups = {}
    for w in raw:
        if len(w) < 8 or not str(w[4]).strip():
            continue
        key = (w[5], w[6])
        groups.setdefault(key, []).append(w)
    lines = []
    for key in sorted(groups):
        ws = sorted(groups[key], key=lambda w: w[7])
        x0 = min(w[0] for w in ws)
        y0 = min(w[1] for w in ws)
        x1 = max(w[2] for w in ws)
        y1 = max(w[3] for w in ws)
        text = " ".join(str(w[4]) for w in ws).strip()
        if not text:
            continue
        lines.append({
            "x": round(x0 / pw, 5), "y": round(y0 / ph, 5),
            "w": round((x1 - x0) / pw, 5), "h": round((y1 - y0) / ph, 5),
            "t": text,
        })
    return lines


def render_image_chalk(src_path: str, out_path: str) -> str:
    """Render an uploaded image to a chalk-on-board PNG at out_path."""
    from PIL import Image
    import numpy as np
    img = Image.open(src_path).convert("RGB")
    rgba, _ = _chalk_from_rgb(np.asarray(img))
    _save_chalk_png(rgba, out_path)
    return out_path


def warm_render():
    """Pre-import/JIT the render stack so the first slide isn't a ~5s cold start.

    fitz + PIL + numpy import (and the first get_pixmap / PNG encode) cost a few
    seconds the first time; doing it once at startup makes every real render fast.
    """
    try:
        import io as _io
        import fitz
        import numpy as np
        from PIL import Image
        doc = fitz.open()
        doc.new_page(width=40, height=40)
        pix = doc.load_page(0).get_pixmap(dpi=72, alpha=False)
        arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
        rgba, _ = _chalk_from_rgb(arr)
        Image.fromarray(rgba, "RGBA").save(_io.BytesIO(), format="PNG", compress_level=1)
        doc.close()
    except Exception:
        pass


def slide_plain_text(slide) -> str:
    """Combined text used for grounding/quiz (verbatim text + figure OCR)."""
    if isinstance(slide, str):
        return slide
    parts = [slide.get("text", "")]
    if slide.get("ocr"):
        parts.append(slide["ocr"])
    return "\n".join(p for p in parts if p and p.strip()).strip()


def _first_line(text: str) -> str:
    for line in (text or "").splitlines():
        line = line.strip()
        if len(line) >= 3:
            return line[:60]
    return "Slide"


def _save_media(media_dir: str, base_name: str, data: bytes, ext: str):
    """Write media bytes to media_dir; return the bare filename or None."""
    if not data:
        return None
    ext = ext if str(ext).startswith(".") else "." + str(ext)
    fn = f"{base_name}{ext.lower()}"
    try:
        with open(os.path.join(media_dir, fn), "wb") as f:
            f.write(data)
        return fn
    except Exception:
        return None


def _dedupe_links(links: list) -> list:
    seen, out = set(), []
    for l in links or []:
        u = (l or {}).get("url")
        if u and u not in seen:
            seen.add(u)
            out.append(l)
    return out


def _text_rich(text: str) -> list:
    return [
        {"blocks": [{"type": "text", "text": t}], "text": t, "ocr": "",
         "links": [], "examinable": True}
        for t in _autosplit(text)
    ]


def _pptx_rich(file_path: str, media_dir: str) -> list:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return _text_rich(_extract_pptx(file_path))
    from PIL import Image

    prs = Presentation(file_path)
    slide_area = float(prs.slide_width or 1) * float(prs.slide_height or 1) or 1.0
    out = []
    for si, slide in enumerate(prs.slides):
        items = []      # (top, left, block)
        ocr_all = []
        links = []
        for shi, shape in enumerate(slide.shapes):
            top = shape.top if shape.top is not None else 0
            left = shape.left if shape.left is not None else 0
            stype = getattr(shape, "shape_type", None)

            # verbatim text + run hyperlinks
            txt = ""
            try:
                if shape.has_text_frame:
                    txt = shape.text_frame.text or ""
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                addr = run.hyperlink.address
                                if addr:
                                    links.append({"text": (run.text or addr).strip(), "url": addr})
                            except Exception:
                                pass
                elif hasattr(shape, "text"):
                    txt = shape.text or ""
            except Exception:
                txt = getattr(shape, "text", "") or ""
            if txt and txt.strip():
                items.append((top, left, {"type": "text", "text": txt.strip()}))

            # pictures (figures) — full-slide scans become text via OCR
            if stype == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    iext = shape.image.ext or "png"
                    frac = 0.0
                    if shape.width and shape.height:
                        frac = (float(shape.width) * float(shape.height)) / slide_area
                    ocr = _ocr_image_obj(Image.open(io.BytesIO(blob))) or ""
                    if frac >= _FULLSLIDE_IMG_FRAC:
                        if ocr:
                            items.append((top, left, {"type": "text", "text": ocr}))
                    else:
                        if ocr:
                            ocr_all.append(ocr)
                        fn = _save_media(media_dir, f"s{si}_p{shi}", blob, iext)
                        if fn:
                            items.append((top, left, {"type": "image", "src": fn, "alt": ocr[:200]}))
                except Exception:
                    pass

        # embedded videos via slide relationships (appended after the text)
        try:
            for rel in slide.part.rels.values():
                rt = (getattr(rel, "reltype", "") or "").lower()
                if "video" in rt or ("media" in rt and "image" not in rt):
                    try:
                        part = rel.target_part
                        pn = str(getattr(part, "partname", "media.mp4"))
                        vext = os.path.splitext(pn)[1] or ".mp4"
                        fn = _save_media(media_dir, f"s{si}_v{abs(hash(pn)) % 100000}", part.blob, vext)
                        if fn:
                            items.append((10 ** 9, 0, {"type": "video", "src": fn}))
                    except Exception:
                        pass
        except Exception:
            pass

        items.sort(key=lambda t: (t[0], t[1]))
        blocks = [b for _, _, b in items]
        text_only = "\n".join(b["text"] for b in blocks if b["type"] == "text")
        out.append({
            "blocks": blocks,
            "text": text_only,
            "ocr": "\n".join(ocr_all),
            "links": _dedupe_links(links),
            "examinable": True,
        })
    return out


def _pdf_rich(file_path: str, media_dir: str) -> list:
    try:
        import fitz  # PyMuPDF
    except Exception:
        return _text_rich(_extract_pdf(file_path))
    from PIL import Image

    doc = fitz.open(file_path)
    out = []
    try:
        for pi in range(doc.page_count):
            page = doc.load_page(pi)
            page_area = float(abs(page.rect.width * page.rect.height)) or 1.0
            native_text = (page.get_text() or "").strip()
            try:
                tblocks = page.get_text("blocks") or []
            except Exception:
                tblocks = []

            img_infos = []  # (xref, y0, area)
            try:
                for img in page.get_images(full=True):
                    xref = img[0]
                    try:
                        rects = page.get_image_rects(xref)
                    except Exception:
                        rects = []
                    if rects:
                        r = rects[0]
                        img_infos.append((xref, float(r.y0), float(abs(r.width * r.height))))
                    else:
                        img_infos.append((xref, 0.0, 0.0))
            except Exception:
                img_infos = []

            items = []  # (y0, block)
            ocr_all = []
            big_img = any((a / page_area) >= _FULLSLIDE_IMG_FRAC for _, _, a in img_infos)

            if big_img and len(native_text) < _MIN_NATIVE_CHARS:
                # scanned page → OCR the whole page as text, show no figure
                ocr = _ocr_pdf_page(doc, pi) or ""
                if ocr:
                    items.append((0.0, {"type": "text", "text": ocr}))
            else:
                for b in tblocks:
                    if len(b) >= 5 and isinstance(b[4], str) and b[4].strip():
                        items.append((float(b[1]), {"type": "text", "text": b[4].strip()}))
                for xref, y0, area in img_infos:
                    if (area / page_area) >= _FULLSLIDE_IMG_FRAC:
                        continue  # page-background scan, not a figure
                    try:
                        info = doc.extract_image(xref)
                        data = info.get("image")
                        iext = info.get("ext", "png")
                        ocr = _ocr_image_obj(Image.open(io.BytesIO(data))) or ""
                        if ocr:
                            ocr_all.append(ocr)
                        fn = _save_media(media_dir, f"p{pi}_i{xref}", data, iext)
                        if fn:
                            items.append((y0, {"type": "image", "src": fn, "alt": ocr[:200]}))
                    except Exception:
                        pass
                if not any(it[1]["type"] == "text" for it in items):
                    ocr = _ocr_pdf_page(doc, pi) or ""
                    if ocr:
                        items.append((0.0, {"type": "text", "text": ocr}))

            links = []
            try:
                for ln in page.get_links():
                    uri = ln.get("uri")
                    if uri:
                        links.append({"text": uri, "url": uri})
            except Exception:
                pass

            items.sort(key=lambda t: t[0])
            blocks = [b for _, b in items]
            text_only = "\n".join(b["text"] for b in blocks if b["type"] == "text")
            out.append({
                "blocks": blocks,
                "text": text_only,
                "ocr": "\n".join(ocr_all),
                "links": _dedupe_links(links),
                "examinable": True,
            })
    finally:
        try:
            doc.close()
        except Exception:
            pass
    return out


def _image_rich(file_path: str, media_dir: str) -> list:
    ocr = _ocr_image_obj(file_path) or ""
    blocks = []
    fn = None
    try:
        with open(file_path, "rb") as f:
            data = f.read()
        iext = os.path.splitext(file_path)[1].lstrip(".") or "png"
        fn = _save_media(media_dir, "img0", data, iext)
    except Exception:
        fn = None
    if ocr:
        blocks.append({"type": "text", "text": ocr})
    if fn:
        blocks.append({"type": "image", "src": fn, "alt": ocr[:200]})
    return [{"blocks": blocks, "text": ocr, "ocr": ocr, "links": [], "examinable": True}]
